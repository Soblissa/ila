from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

src='docs/praesentationen/ILA_Pilot_Design-Entlastung_V4_HMKB.pptx'
out='docs/praesentationen/ILA_Pilot_Design-Entlastung_V4_HMKB_mit_Symbolfolie.pptx'

prs=Presentation(src)
layout=prs.slide_layouts[6] if len(prs.slide_layouts)>6 else prs.slide_layouts[-1]
slide=prs.slides.add_slide(layout)

# title
left,top,width,height=Inches(0.8),Inches(0.5),Inches(12),Inches(0.8)
box=slide.shapes.add_textbox(left,top,width,height)
tf=box.text_frame
p=tf.paragraphs[0]
p.text='Symbolübersicht für die 6 Inhaltsfolien'
p.font.size=Pt(34)
p.font.bold=True
p.font.color.rgb=RGBColor(0x0B,0x4F,0x9C)

items=[
    ('◉','Ausgangslage'),
    ('✦','Vorschlag'),
    ('⚙','Pilot-Setup'),
    ('⏱','Erfolgsmessung'),
    ('⚠','Pilotrisiken'),
    ('◎','Pilotziele'),
]

start_y=1.6
for i,(sym,label) in enumerate(items):
    y=Inches(start_y + i*0.75)
    b=slide.shapes.add_textbox(Inches(1.2),y,Inches(10),Inches(0.5))
    tf=b.text_frame
    p=tf.paragraphs[0]
    p.text=f'{sym}  {label}'
    p.font.size=Pt(30)
    p.font.color.rgb=RGBColor(0x0B,0x4F,0x9C)

prs.save(out)
print(out)

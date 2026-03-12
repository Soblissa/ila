from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

out='docs/praesentationen/ILA_Piktogramme_nur_Symbole.pptx'
prs=Presentation('PowerPoint-Vorlage_HMKB.pptx')
layout=prs.slide_layouts[7] if len(prs.slide_layouts)>7 else prs.slide_layouts[-1]
slide=prs.slides.add_slide(layout)

blue=RGBColor(0x0B,0x4F,0x9C)
# Nur Symbole, keine Überschrift, keine Labels
symbols=['⌕','✦','⚙','⏱','⚠','◎']
# 3x2 Raster
positions=[
    (1.8,1.6),(5.6,1.6),(9.4,1.6),
    (1.8,4.2),(5.6,4.2),(9.4,4.2),
]
for sym,(x,y) in zip(symbols,positions):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(2.2),Inches(1.4))
    tf=box.text_frame
    p=tf.paragraphs[0]
    p.text=sym
    p.font.size=Pt(88)
    p.font.bold=True
    p.font.color.rgb=blue

prs.save(out)
print(out)

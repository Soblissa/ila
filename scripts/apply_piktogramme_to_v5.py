from pptx import Presentation

src='docs/praesentationen/ILA_Pilot_Design-Entlastung_V4_HMKB.pptx'
out='docs/Präsentation/ILA_Pilot_Design-Entlastung_V5_HMKB.pptx'

prs=Presentation(src)

# Map slide index (0-based) to new title with piktogramm
updates = {
    1: '🔍 Ausgangslage',
    2: '💡 Vorschlag',
    3: '⚙️ Pilot-Setup (4 Wochen)',
    4: '📊 Erfolgsmessung (im Agententeam)',
    5: '⚠️ Pilotrisiken',
    6: '🎯 Pilotziele',
}

for idx, title in updates.items():
    if idx < len(prs.slides):
        slide = prs.slides[idx]
        if slide.shapes.title is not None:
            slide.shapes.title.text = title

prs.save(out)
print(out)

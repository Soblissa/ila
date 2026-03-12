from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

out='docs/Präsentation/ILA_Piktogramme_Uebersichtenfolie.pptx'
prs=Presentation('PowerPoint-Vorlage_HMKB.pptx')
layout=prs.slide_layouts[7] if len(prs.slide_layouts)>7 else prs.slide_layouts[-1]
slide=prs.slides.add_slide(layout)

blue=RGBColor(0x0B,0x4F,0x9C)

# Title
box=slide.shapes.add_textbox(Inches(0.8),Inches(0.5),Inches(12),Inches(0.8))
p=box.text_frame.paragraphs[0]
p.text='Piktogramm-Vorschläge für Inhaltsfolien'
p.font.size=Pt(34)
p.font.bold=True
p.font.color.rgb=blue

items=[
    ('🔍','Ausgangslage'),
    ('💡','Vorschlag'),
    ('⚙️','Pilot-Setup'),
    ('📊','Erfolgsmessung'),
    ('⚠️','Pilotrisiken'),
    ('🎯','Pilotziele'),
]

start_y=1.6
for i,(icon,label) in enumerate(items):
    y=Inches(start_y+i*0.8)
    b=slide.shapes.add_textbox(Inches(1.0),y,Inches(11),Inches(0.6))
    tf=b.text_frame
    p=tf.paragraphs[0]
    p.text=f'{icon}  {label}'
    p.font.size=Pt(30)
    p.font.color.rgb=blue

prs.save(out)
print(out)

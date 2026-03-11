from pptx import Presentation
from pptx.util import Pt

TEMPLATE = 'PowerPoint-Vorlage_HMKB.pptx'
OUT = 'docs/praesentationen/ILA_Pilot_Design-Entlastung_2026-03-14.pptx'

slides_data = [
    ("ILA schneller machen: Design-Entlastung durch Agenten-Pilot", [
        "Kontrollierter Test statt Grundsatzdebatte",
        "Ziel heute: Freigabe für einen 4-Wochen-Pilot",
        "Design bei Routinearbeit entlasten, Qualität sichern"
    ]),
    ("Ausgangslage", [
        "Hohe Komplexität im Projektumfeld",
        "Reibung zwischen Fachlichkeit und Design",
        "Iterationen dauern zu lange"
    ]),
    ("Was der Agent konkret übernimmt", [
        "Figma-Erstentwürfe auf Basis des bestehenden Designs",
        "Schnelle Varianten für wiederkehrende Screen-Typen",
        "Mensch bleibt in fachlicher Entscheidung und Freigabe"
    ]),
    ("Pilot-Setup in 4 Wochen", [
        "Scope: 1 bis 2 Seitentypen mit hoher Wiederholbarkeit",
        "Ablauf: Anforderung -> Agent-Entwurf -> Figma -> Review",
        "Klare Rollen: Fachseite, Design, Technik, Projektleitung"
    ]),
    ("Nutzen und Messung", [
        "Durchlaufzeit pro Screen (vorher/nachher)",
        "Anzahl Korrekturschleifen",
        "Konsistenz in Standard-Screens",
        "Entlastung des Design-Engpasses"
    ]),
    ("Risiken und Absicherung", [
        "Uneinheitliche Ergebnisse -> Regelset + Referenzscreens vor Start",
        "Rollenunsicherheit -> klare Verantwortlichkeiten",
        "Grundsatzdebatte -> Pilotfokus + messbare Kriterien"
    ]),
    ("Nicht-Ziele", [
        "Kein vollständiger Ersatz von Design-Rollen",
        "Kein Big-Bang-Umbau",
        "Keine ungeprüfte Automatisierung"
    ]),
    ("Entscheidung heute", [
        "4-Wochen-Pilot freigeben",
        "1 bis 2 Seitentypen festlegen",
        "Verantwortliche benennen",
        "Go/No-Go Termin am Pilotende setzen"
    ]),
]

prs = Presentation(TEMPLATE)

# Prefer title+content layout when available
layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

for title, bullets in slides_data:
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    else:
        tx = slide.shapes.add_textbox(Pt(40), Pt(20), Pt(860), Pt(60)).text_frame
        tx.text = title

    # Body placeholder fallback
    body = None
    for shape in slide.placeholders:
        if shape.is_placeholder and shape.placeholder_format.idx != 0:
            if hasattr(shape, 'text_frame'):
                body = shape.text_frame
                break

    if body is None:
        body = slide.shapes.add_textbox(Pt(50), Pt(120), Pt(850), Pt(380)).text_frame

    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(24)

prs.save(OUT)
print(OUT)

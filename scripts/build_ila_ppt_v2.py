from pptx import Presentation

TEMPLATE='PowerPoint-Vorlage_HMKB.pptx'
OUT='docs/praesentationen/ILA_Pilot_Design-Entlastung_V2_HMKB.pptx'

# layout indices from template
TITLE=0
CONTENT=3
TWOCOL=5

slides=[
    (TITLE, 'ILA schneller machen', ['Pilot zur Design-Entlastung mit Agenten','Freitag: Entscheidungsgrundlage']),
    (CONTENT, 'Ausgangslage', ['Komplexes Umfeld mit hoher Abstimmungsdichte','Reibung zwischen Fachlichkeit und Design','Standard-Screens binden viel Designkapazität']),
    (CONTENT, 'Vorschlag', ['Agent erstellt Figma-Erstentwürfe auf Basis bestehender Muster','Mensch bleibt für Entscheidung und Freigabe verantwortlich','Fokus auf wiederkehrende Seitentypen']),
    (CONTENT, 'Pilot-Setup (4 Wochen)', ['Scope: 1 bis 2 Seitentypen','Ablauf: Anforderung -> Entwurf -> Review -> Freigabe','Klare Rollen für Fachseite, Design, Technik, Projektleitung']),
    (CONTENT, 'Erfolgsmessung', ['Durchlaufzeit pro Screen','Anzahl Korrekturschleifen','Konsistenz nach Checkliste']),
    (CONTENT, 'Risiken und Absicherung', ['Uneinheitliche Ergebnisse -> Regelset + Referenzscreens','Rollenunklarheit -> Verantwortlichkeiten vorab klären','Grundsatzdebatte -> Pilot eng auf Nutzen begrenzen']),
    (CONTENT, 'Nicht-Ziele', ['Kein vollständiger Ersatz von Design-Rollen','Kein Big-Bang-Umbau','Keine ungeprüfte Automatisierung']),
    (CONTENT, 'Entscheidung heute', ['Pilot für 4 Wochen freigeben','Seitentypen und Verantwortliche festlegen','Go/No-Go Termin am Pilotende setzen']),
]

prs=Presentation(TEMPLATE)

for layout_idx,title,bullets in slides:
    slide=prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # title
    if slide.shapes.title is not None:
        slide.shapes.title.text=title
    else:
        # fallback first placeholder with text_frame
        for shp in slide.placeholders:
            if hasattr(shp,'text_frame'):
                shp.text_frame.text=title
                break

    # content placeholder
    body=None
    # prefer non-title placeholder idx 1
    for ph in slide.placeholders:
        if hasattr(ph,'text_frame') and ph != slide.shapes.title:
            body=ph.text_frame
            break

    if body is not None and bullets:
        body.clear()
        for i,b in enumerate(bullets):
            p = body.paragraphs[0] if i==0 else body.add_paragraph()
            p.text=b
            p.level=0

prs.save(OUT)
print(OUT)

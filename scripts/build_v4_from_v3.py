from pptx import Presentation

SRC='ILA_Pilot_Design-Entlastung_V3_HMKB.pptx'
OUT='ILA_Pilot_Design-Entlastung_V4_HMKB.pptx'

prs=Presentation(SRC)

# Helpers

def set_title(slide, text):
    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        return
    for sh in slide.shapes:
        if hasattr(sh, 'text_frame'):
            sh.text_frame.text = text
            return

def get_body_tf(slide):
    # non-title placeholder with text_frame
    for sh in slide.shapes:
        if sh == slide.shapes.title:
            continue
        if hasattr(sh, 'text_frame'):
            return sh.text_frame
    return None

def set_bullets(slide, bullets):
    tf=get_body_tf(slide)
    if tf is None:
        return
    tf.clear()
    for i,b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=b
        p.level=0

# Slide 1
s=prs.slides[0]
# keep existing title style, set both prominent text boxes safely
# title + subtitle in first non-title textframe
set_title(s, 'ila schneller machen')
set_bullets(s, ['Von Reibung zu Wirkung: Pilot zur Beschleunigung des Designs'])

# Slide 2
s=prs.slides[1]
set_title(s, 'Ausgangslage')
set_bullets(s, [
    'Hohe Komplexität und Reibung zwischen Fachlichkeit und Design',
    'Wiederkehrende Standardscreens (z. B. Formular-, Listen- und Detailansichten) binden viel Designkapazität',
    'Iterationen dauern zu lange',
    'Angespannte Ressourcenlage bei hohem Stakeholder-Druck',
    'Design ist aktueller Flaschenhals der Produktentwicklung',
])

# Slide 3
s=prs.slides[2]
set_title(s, 'Vorschlag')
set_bullets(s, [
    'Pilotversuch in kontrolliertem technischem Setup.',
    'Klare Governance für Datenschutz, Rollen und Freigaben.',
    'Agent erstellt Figma-Erstentwürfe auf Basis definierter Designspezifikationen.',
    'Fokus auf aktuelles Teilziel: BSLRR (besondere Schwierigkeiten im Lesen, Rechtschreiben und Rechnen).',
    'Mensch bleibt für Entscheidung und Freigabe verantwortlich.',
])

# Slide 4
s=prs.slides[3]
set_title(s, 'Pilot-Setup (4 Wochen)')
set_bullets(s, [
    'Voraussetzung: Definierte Designspezifikationen und Referenzscreens liegen vor.',
    'Setup: Kontrollierter technischer Rahmen mit klarer Governance (Datenschutz, Rollen, Freigaben und Verantwortlichkeiten).',
    'Scope: 1 bis 2 Seitentypen mit Fokus auf BSLRR.',
    'Rollen: Fachseite (Sebastian Felgenhauer), Technik (Markus Deibel) und Projektleitung (Sarah Bonhoff).',
    'Ablauf: Fachliche Anforderung aufnehmen und schärfen -> Figma-Erstentwurf auf Basis der Designspezifikation erstellen -> Review mit dem gesamten ILA-Team.',
])

# Slide 5
s=prs.slides[4]
set_title(s, 'Erfolgsmessung (im Agententeam)')
set_bullets(s, [
    'Agentenzeit pro Entwurf (reine Lauf-/Erstellzeit des Agenten).',
    'Time-to-First-Draft (Anforderung bis erster belastbarer Figma-Entwurf).',
    'End-to-End-Zeit pro Screen (Anforderung bis reviewfähiger Entwurf).',
    'Rework-Aufwand (Anzahl und Dauer der Korrekturschleifen).',
])

# Slide 6
s=prs.slides[5]
set_title(s, 'Pilotrisiken')
set_bullets(s, [
    'Risiko 1: Der Test zeigt, dass Agenteneinsatz im aktuellen Rahmen nicht realisierbar ist.',
    'Mögliche Ursachen: unklare Designspezifikation, zu viele Sonderfälle, instabile Mitwirkung im Pilotprozess.',
    'Risiko 2: Arbeitszeit von Sebastian Felgenhauer, Markus Deibel und Sarah Bonhoff wird gebunden, ohne verwertbaren Mehrwert.',
])

# Slide 7
s=prs.slides[6]
set_title(s, 'Pilotziele')
set_bullets(s, [
    'Machbarkeitsprüfung: Ist Agenteneinsatz als Ergänzung zur Designrolle im Projekt tragfähig?',
    'Evaluierung: Welchen konkreten Nutzen bringt der Ansatz im Pilotbetrieb?',
    'Skalierungspotential: In welchem Umfang ist ein Rollout sinnvoll und realistisch?',
    'Bereitstellung einer Entscheidungsvorlage: Klare Grundlage für die weitere Entscheidung nach Pilotabschluss.',
])

# Remove slide 8 (index 7)
slide_id = prs.slides._sldIdLst[7]
prs.slides._sldIdLst.remove(slide_id)

prs.save(OUT)
print(OUT)
